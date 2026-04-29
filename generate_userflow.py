#!/usr/bin/env python3
"""Generate UCInsure user flow diagram as a native PPTX file."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE


def rgb(r, g, b):
    return RGBColor(r, g, b)


C_START   = rgb(0x2E, 0x86, 0xAB)
C_PROC    = rgb(0xEE, 0xEE, 0xEE)
C_DIAMOND = rgb(0xFF, 0xF0, 0xCC)
C_OUTPUT  = rgb(0xD0, 0xF5, 0xE4)
C_ERROR   = rgb(0xFF, 0xD9, 0xD9)
C_TRAIN   = rgb(0xE8, 0xD5, 0xF5)
C_LINE    = rgb(0x55, 0x55, 0x55)
C_TEXT    = rgb(0x22, 0x22, 0x22)

ROUNDED = 5
DIAMOND = 4
OVAL    = 9


def node(slide, text, cx, cy, w, h,
         shape_id=ROUNDED, fill=C_PROC, fsize=10, bold=False):
    s = slide.shapes.add_shape(
        shape_id,
        Inches(cx - w / 2), Inches(cy - h / 2),
        Inches(w), Inches(h),
    )
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = C_LINE
    s.line.width = Pt(1.0)
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fsize)
    run.font.bold = bold
    run.font.color.rgb = C_TEXT
    return s


def arrow(slide, x1, y1, x2, y2, dashed=False):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2),
    )
    c.line.color.rgb = C_LINE
    c.line.width = Pt(1.2)
    if dashed:
        from pptx.oxml.ns import qn
        import lxml.etree as etree
        prstDash = etree.SubElement(c.line._ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')
    return c


def lbl(slide, text, cx, cy):
    tb = slide.shapes.add_textbox(
        Inches(cx - 0.4), Inches(cy - 0.14),
        Inches(0.8), Inches(0.28),
    )
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.color.rgb = C_LINE


def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)   # widescreen 16:9
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── Layout ──────────────────────────────────────────────────────────────
    MCX = 6.67   # centre of slide
    RX  = 11.2   # right branch
    LX  = 2.8    # left branch
    NW  = 3.2    # normal node width
    IW  = 2.0    # input node width
    SW  = 2.0    # side node width
    OW  = 2.2    # output node width
    NH  = 0.46   # node height
    DW  = 1.65   # diamond width
    DH  = 0.58   # diamond height
    GAP = 0.65   # row pitch — 11 rows × 0.65 = 7.15 fits within 7.5"

    Y = [0.38 + i * GAP for i in range(11)]
    # Y[10] = 0.38 + 10*0.65 = 6.88,  bottom = 7.11 ✓

    def t(cy, h=NH): return cy - h / 2
    def b(cy, h=NH): return cy + h / 2

    # ── Row 0: START ────────────────────────────────────────────────────────
    node(slide, "🚀  Start", MCX, Y[0], 1.9, 0.44,
         shape_id=OVAL, fill=C_START, fsize=12, bold=True)

    # ── Row 1: Entry decision ───────────────────────────────────────────────
    node(slide, "How to\nProvide Data?", MCX, Y[1], DW, DH,
         DIAMOND, C_DIAMOND, 10)
    arrow(slide, MCX, b(Y[0], 0.44), MCX, t(Y[1], DH))

    # ── Row 2: Input options (3 nodes) ──────────────────────────────────────
    IXs = [MCX - 2.8, MCX, MCX + 2.8]
    for ix, txt in zip(IXs, ["Upload\nCSV File", "Manual\nEntry", "Refresh\nDataset"]):
        node(slide, txt, ix, Y[2], IW, NH, fill=C_PROC, fsize=10)
    arrow(slide, MCX - DW/2, Y[1], IXs[0] + IW/2, t(Y[2]))
    arrow(slide, MCX,        b(Y[1], DH), MCX, t(Y[2]))
    arrow(slide, MCX + DW/2, Y[1], IXs[2] - IW/2, t(Y[2]))

    # ── Row 3: Validate ─────────────────────────────────────────────────────
    node(slide, "Data\nValid?", MCX, Y[3], DW, DH, DIAMOND, C_DIAMOND, 10)
    for ix in IXs:
        arrow(slide, ix, b(Y[2]), MCX, t(Y[3], DH))

    # Error branch (right)
    node(slide, "⚠  Validation Error\nReview & Fix", RX, Y[3], SW, NH,
         fill=C_ERROR, fsize=10)
    arrow(slide, MCX + DW/2, Y[3], RX - SW/2, Y[3])
    lbl(slide, "No", (MCX + DW/2 + RX - SW/2) / 2, Y[3] - 0.17)
    arrow(slide, RX, t(Y[3]), RX, Y[1], dashed=True)
    arrow(slide, RX, Y[1], MCX + DW/2, Y[1], dashed=True)
    lbl(slide, "retry", RX + 0.35, (Y[1] + Y[3]) / 2)

    # ── Row 4: Preprocess ───────────────────────────────────────────────────
    node(slide, "Preprocess & Feature Engineering", MCX, Y[4], NW, NH,
         fill=C_PROC, fsize=10)
    arrow(slide, MCX, b(Y[3], DH), MCX, t(Y[4]))
    lbl(slide, "Yes", MCX + 0.3, (b(Y[3], DH) + t(Y[4])) / 2)

    # ── Row 5: Train? ───────────────────────────────────────────────────────
    node(slide, "Train\nModels?", MCX, Y[5], DW, DH, DIAMOND, C_DIAMOND, 10)
    arrow(slide, MCX, b(Y[4]), MCX, t(Y[5], DH))

    node(slide, "🔚  Exit", RX, Y[5], 1.6, 0.44,
         shape_id=OVAL, fill=C_ERROR, fsize=10)
    arrow(slide, MCX + DW/2, Y[5], RX - 0.8, Y[5])
    lbl(slide, "No", (MCX + DW/2 + RX - 0.8) / 2, Y[5] - 0.17)

    # ── Row 6: Train models ─────────────────────────────────────────────────
    node(slide, "Run Model Training\nRandom Forest  ·  Mean Score  ·  Frequency",
         MCX, Y[6], NW, NH, fill=C_TRAIN, fsize=10)
    arrow(slide, MCX, b(Y[5], DH), MCX, t(Y[6]))
    lbl(slide, "Yes", MCX + 0.3, (b(Y[5], DH) + t(Y[6])) / 2)

    # ── Row 7: Metrics review? ──────────────────────────────────────────────
    node(slide, "Metrics\nReview?", MCX, Y[7], DW, DH, DIAMOND, C_DIAMOND, 10)
    arrow(slide, MCX, b(Y[6]), MCX, t(Y[7], DH))

    node(slide, "Best Model\nAuto-selected", LX, Y[7], SW, NH, fill=C_PROC, fsize=10)
    arrow(slide, MCX - DW/2, Y[7], LX + SW/2, Y[7])
    lbl(slide, "Auto", (MCX - DW/2 + LX + SW/2) / 2, Y[7] - 0.17)

    node(slide, "View & Select\nBest Metrics", RX, Y[7], SW, NH, fill=C_PROC, fsize=10)
    arrow(slide, MCX + DW/2, Y[7], RX - SW/2, Y[7])
    lbl(slide, "Manual", (MCX + DW/2 + RX - SW/2) / 2, Y[7] - 0.17)

    arrow(slide, LX, b(Y[7]), MCX - 0.4, t(Y[8]))
    arrow(slide, RX, b(Y[7]), MCX + 0.4, t(Y[8]))

    # ── Row 8: Predict ──────────────────────────────────────────────────────
    node(slide, "Enter Input & Run Prediction", MCX, Y[8], NW, NH,
         fill=C_PROC, fsize=10)

    # ── Row 9: Result ───────────────────────────────────────────────────────
    node(slide, "🎯  Risk Level:  Low  ·  Medium  ·  High",
         MCX, Y[9], NW + 0.6, NH, fill=C_OUTPUT, fsize=12, bold=True)
    arrow(slide, MCX, b(Y[8]), MCX, t(Y[9]))

    # ── Row 10: Output options via horizontal bus ────────────────────────────
    OXs   = [MCX - 4.6, MCX - 1.5, MCX + 1.5, MCX + 4.6]
    OTxts = ["Prediction\nExplanation", "Scenario\nSimulation",
              "Comparison\nReport", "✅  Done"]
    OFills = [C_OUTPUT, C_OUTPUT, C_OUTPUT, C_START]
    for ox, ot, of in zip(OXs, OTxts, OFills):
        node(slide, ot, ox, Y[10], OW, NH, fill=of, fsize=10)

    # Bus: vertical stem from result centre down to bus Y, then horizontal
    # spanning all outputs, then vertical drops to each output box
    BUS_Y = b(Y[9]) + (t(Y[10]) - b(Y[9])) * 0.5   # halfway between rows 9 and 10
    arrow(slide, MCX, b(Y[9]), MCX, BUS_Y)           # stem down from result
    arrow(slide, OXs[0], BUS_Y, OXs[-1], BUS_Y)      # horizontal bus
    for ox in OXs:
        arrow(slide, ox, BUS_Y, ox, t(Y[10]))         # drop to each output

    # New prediction loop back to Predict (dashed, left margin)
    LOOP_X = OXs[0] - OW/2 - 0.15
    arrow(slide, OXs[0] - OW/2, Y[10], LOOP_X, Y[10], dashed=True)
    arrow(slide, LOOP_X, Y[10], LOOP_X, Y[8], dashed=True)
    arrow(slide, LOOP_X, Y[8], MCX - NW/2, Y[8], dashed=True)
    lbl(slide, "↺ new", LOOP_X - 0.45, (Y[8] + Y[10]) / 2)

    out = "ucinsure_userflow.pptx"
    prs.save(out)
    print(f"Saved → {out}")


if __name__ == "__main__":
    main()
