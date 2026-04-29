#!/usr/bin/env python3
"""Generate UCInsure system architecture diagram as a native PPTX file."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE


def rgb(r, g, b):
    return RGBColor(r, g, b)


# Colour palette
C_ACTOR   = rgb(0xD6, 0xE8, 0xFF)  # blue  – actor
C_PROC    = rgb(0xEE, 0xEE, 0xEE)  # grey  – process
C_DIAMOND = rgb(0xFF, 0xF0, 0xCC)  # amber – decision
C_OUTPUT  = rgb(0xD0, 0xF5, 0xE4)  # green – output
C_ERROR   = rgb(0xFF, 0xD9, 0xD9)  # red   – error
C_LINE    = rgb(0x55, 0x55, 0x55)
C_TEXT    = rgb(0x22, 0x22, 0x22)

ROUNDED_RECT = 5
DIAMOND      = 4


def node(slide, text, cx, cy, w, h,
         shape_id=ROUNDED_RECT, fill=C_PROC, fsize=12, bold=False):
    s = slide.shapes.add_shape(
        shape_id,
        Inches(cx - w / 2), Inches(cy - h / 2),
        Inches(w), Inches(h),
    )
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = C_LINE
    s.line.width = Pt(1.0)
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fsize)
    run.font.bold = bold
    run.font.color.rgb = C_TEXT
    return s


def arrow(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2),
    )
    c.line.color.rgb = C_LINE
    c.line.width = Pt(1.3)
    return c


def label(slide, text, cx, cy):
    tb = slide.shapes.add_textbox(
        Inches(cx - 0.38), Inches(cy - 0.14),
        Inches(0.76), Inches(0.28),
    )
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = C_LINE


def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── Layout constants ─────────────────────────────────────────────────────
    MCX = 4.0    # main column x-centre
    SCX = 8.1    # side column x-centre
    NW  = 3.2    # normal node width
    SW  = 2.0    # side node width
    NH  = 0.52   # node height
    DW  = 1.60   # diamond width
    DH  = 0.62   # diamond height
    GAP = 0.80   # row pitch

    Y = [0.55 + i * GAP for i in range(9)]
    # Y[0]=0.55  Y[1]=1.35  Y[2]=2.15  Y[3]=2.95  Y[4]=3.75
    # Y[5]=4.55  Y[6]=5.35  Y[7]=6.15  Y[8]=6.95

    # ── Nodes ────────────────────────────────────────────────────────────────
    node(slide, "👤  User / Analyst",
         MCX, Y[0], NW, NH, fill=C_ACTOR, fsize=13, bold=True)

    node(slide, "Data  Source?",
         MCX, Y[1], DW, DH, shape_id=DIAMOND, fill=C_DIAMOND, fsize=11)

    node(slide, "Data Ingestion\nUpload  ·  Refresh  ·  Manual Entry",
         MCX, Y[2], NW, NH, fsize=12)

    node(slide, "Valid?",
         MCX, Y[3], DW, DH, shape_id=DIAMOND, fill=C_DIAMOND, fsize=11)
    node(slide, "⚠  Validation Error",
         SCX, Y[3], SW, NH, fill=C_ERROR, fsize=11)

    node(slide, "Preprocessing & Feature Engineering",
         MCX, Y[4], NW, NH, fsize=12)

    node(slide, "Model Training\nRandom Forest  ·  Mean Score  ·  Frequency",
         MCX, Y[5], NW, NH, fsize=12)

    node(slide, "Select  Model?",
         MCX, Y[6], DW, DH, shape_id=DIAMOND, fill=C_DIAMOND, fsize=11)
    node(slide, "Metrics Review",
         SCX, Y[6], SW, NH, fsize=11)

    OUT_H = NH + 0.08
    node(slide, "🎯  Risk Prediction\nLow  ·  Medium  ·  High",
         MCX, Y[7], NW, OUT_H, fill=C_OUTPUT, fsize=13, bold=True)

    for text, xc in [("Explanation", 2.3), ("Simulation", 4.0), ("Report", 5.7)]:
        node(slide, text, xc, Y[8], 1.5, NH, fill=C_OUTPUT, fsize=11)

    # ── Edge helpers ─────────────────────────────────────────────────────────
    def top(cy, h=NH): return cy - h / 2
    def bot(cy, h=NH): return cy + h / 2

    # ── Main-flow arrows ──────────────────────────────────────────────────────
    arrow(slide, MCX, bot(Y[0]),        MCX, top(Y[1], DH))
    arrow(slide, MCX, bot(Y[1], DH),    MCX, top(Y[2]))
    arrow(slide, MCX, bot(Y[2]),        MCX, top(Y[3], DH))
    arrow(slide, MCX, bot(Y[3], DH),    MCX, top(Y[4]))
    label(slide, "Yes", MCX + 0.28, (bot(Y[3], DH) + top(Y[4])) / 2)
    arrow(slide, MCX, bot(Y[4]),        MCX, top(Y[5]))
    arrow(slide, MCX, bot(Y[5]),        MCX, top(Y[6], DH))
    arrow(slide, MCX, bot(Y[6], DH),    MCX, top(Y[7], OUT_H))
    label(slide, "Auto", MCX + 0.3, (bot(Y[6], DH) + top(Y[7], OUT_H)) / 2)
    for xc in [2.3, 4.0, 5.7]:
        arrow(slide, xc, bot(Y[7], OUT_H), xc, top(Y[8]))

    # ── Side-branch arrows ────────────────────────────────────────────────────
    # Valid? —No→ Error
    arrow(slide, MCX + DW / 2, Y[3], SCX - SW / 2, Y[3])
    label(slide, "No", (MCX + DW / 2 + SCX - SW / 2) / 2, Y[3] - 0.17)

    # Error —retry→ up then left to User
    arrow(slide, SCX, top(Y[3]), SCX, Y[0])
    arrow(slide, SCX, Y[0], MCX + NW / 2, Y[0])
    label(slide, "retry", SCX + 0.32, (Y[0] + Y[3]) / 2)

    # Select? —Review→ Metrics
    arrow(slide, MCX + DW / 2, Y[6], SCX - SW / 2, Y[6])
    label(slide, "Review", (MCX + DW / 2 + SCX - SW / 2) / 2, Y[6] - 0.17)

    # Metrics —→ Risk Prediction
    arrow(slide, SCX, bot(Y[6]), SCX, Y[7])
    arrow(slide, SCX, Y[7], MCX + NW / 2, Y[7])

    # ── Write-up text box (right side of slide) ───────────────────────────────
    TX = 7.0   # left edge of text column
    TW = 2.8   # text box width
    TY = 0.40  # top of text box
    TH = 6.70  # height

    from pptx.util import Emu
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    tb = slide.shapes.add_textbox(Inches(TX), Inches(TY), Inches(TW), Inches(TH))
    tf = tb.text_frame
    tf.word_wrap = True

    TITLE  = "UCInsure — Risk Prediction Pipeline"
    BULLETS = [
        "Ingests NFIP claims data via upload, manual entry, or refresh",
        "Validates, preprocesses, and engineers features automatically",
        "Trains & selects the best ML model (Random Forest, Mean Score, Frequency)",
        "Outputs Low / Medium / High risk with explanation & simulation",
    ]

    # Title paragraph
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = TITLE
    r0.font.size = Pt(13)
    r0.font.bold = True
    r0.font.color.rgb = C_TEXT

    # Bullet paragraphs
    for bullet in BULLETS:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = f"•  {bullet}"
        r.font.size = Pt(11)
        r.font.color.rgb = C_TEXT

    # ── Save ──────────────────────────────────────────────────────────────────
    out = "ucinsure_architecture.pptx"
    prs.save(out)
    print(f"Saved → {out}")


if __name__ == "__main__":
    main()
